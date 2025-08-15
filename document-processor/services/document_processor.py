import fitz  # PyMuPDF
import pytesseract
from PIL import Image
from docx import Document as DocxDocument
import os
import uuid
import time
from typing import Optional, Tuple, List, Dict
import logging
from pathlib import Path
import io
import json
from datetime import datetime
import re

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Document processing service using PyMuPDF and OCR"""
    
    def __init__(self):
        self.supported_formats = {
            'pdf': self._process_pdf,
            'docx': self._convert_docx_to_pdf,
            'doc': self._convert_doc_to_pdf,
            'xlsx': self._convert_xlsx_to_pdf,
            'xls': self._convert_xls_to_pdf,
            'pptx': self._convert_pptx_to_pdf,
            'ppt': self._convert_ppt_to_pdf,
            'csv': self._convert_csv_to_pdf,
            'png': self._convert_image_to_pdf,
            'jpg': self._convert_image_to_pdf,
            'jpeg': self._convert_image_to_pdf,
        }
    
    def _sanitize_extracted_text(self, text: str) -> str:
        if not text:
            return ""
        
        try:
            if isinstance(text, bytes):
                # Try UTF-8 first, then latin-1 as fallback
                try:
                    text = text.decode('utf-8')
                except UnicodeDecodeError:
                    text = text.decode('latin-1', errors='ignore')
            
            # Remove or replace UTF-8 replacement characters
            text = text.replace('\ufffd', '')  # Remove replacement characters
            
            # Remove null characters and other control characters that can break JSON
            text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)
            
            # Remove additional problematic Unicode control characters
            text = re.sub(r'[\u0080-\u009F]', '', text)
            
            # Normalize whitespace
            text = re.sub(r'\s+', ' ', text.strip())
            
            # Ensure the result is valid UTF-8
            encoded = text.encode('utf-8', errors='ignore')
            sanitized = encoded.decode('utf-8')
            
            logger.debug(f"Text sanitization: {len(text)} -> {len(sanitized)} characters")
            return sanitized
            
        except Exception as e:
            logger.error(f"Error sanitizing text: {str(e)}")
            # Return a safe fallback
            return re.sub(r'[^\x20-\x7E\s]', '', str(text))
    
    async def process_document(self, file_path: str, original_filename: str) -> Tuple[str, str, float]:
        """
        Process document: convert to PDF if needed and extract text using OCR
        
        Args:
            file_path: Path to the uploaded file
            original_filename: Original filename with extension
            
        Returns:
            Tuple of (pdf_path, extracted_text, processing_time)
        """
        start_time = time.time()
        
        try:
            # Get file extension
            file_ext = Path(original_filename).suffix.lower().lstrip('.')
            
            if file_ext not in self.supported_formats:
                raise ValueError(f"Unsupported file format: {file_ext}")
            
            # Convert to PDF if needed
            if file_ext == 'pdf':
                pdf_path = file_path
            else:
                pdf_path = await self.supported_formats[file_ext](file_path)
            
            # Extract text using PyMuPDF + OCR
            extracted_text = await self._extract_text_with_ocr(pdf_path)
            
            processing_time = time.time() - start_time
            
            return pdf_path, extracted_text, processing_time
            
        except Exception as e:
            logger.error(f"Error processing document {original_filename}: {str(e)}")
            raise
    
    async def _process_pdf(self, file_path: str) -> str:
        """Process existing PDF file"""
        return file_path
    
    async def _convert_docx_to_pdf(self, file_path: str) -> str:
        """Convert DOCX to PDF using python-docx and PyMuPDF"""
        try:
            # Extract text from DOCX
            docx_document = DocxDocument(file_path)
            text_content = []
            
            for paragraph in docx_document.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            full_text = '\n'.join(text_content)
            
            # Create PDF from text using PyMuPDF
            pdf_path = file_path.replace('.docx', '.pdf')
            pdf_doc = fitz.open()  # Create new PDF
            page = pdf_doc.new_page()
            
            # Insert text with basic formatting
            text_rect = fitz.Rect(50, 50, 550, 750)
            page.insert_textbox(text_rect, full_text, fontsize=11, fontname="helv")
            
            pdf_doc.save(pdf_path)
            pdf_doc.close()
            
            return pdf_path
            
        except Exception as e:
            logger.error(f"Error converting DOCX to PDF: {str(e)}")
            raise
    
    async def _convert_doc_to_pdf(self, file_path: str) -> str:
        """Convert DOC to PDF (basic implementation)"""
        # For .doc files, we'll treat them as binary and try to extract what we can
        # In production, you might want to use python-docx2txt or similar
        pdf_path = file_path.replace('.doc', '.pdf')
        
        # Create a simple PDF with a message about DOC conversion
        pdf_doc = fitz.open()
        page = pdf_doc.new_page()
        
        text = "DOC file uploaded. For better text extraction, please upload DOCX or PDF format."
        text_rect = fitz.Rect(50, 50, 550, 750)
        page.insert_textbox(text_rect, text, fontsize=11, fontname="helv")
        
        pdf_doc.save(pdf_path)
        pdf_doc.close()
        
        return pdf_path
    
    async def _convert_xlsx_to_pdf(self, file_path: str) -> str:
        """Convert XLSX to PDF using openpyxl"""
        try:
            from openpyxl import load_workbook
            
            # Load the Excel workbook
            workbook = load_workbook(file_path, data_only=True)
            text_content = []
            
            for sheet_name in workbook.sheetnames:
                worksheet = workbook[sheet_name]
                text_content.append(f"=== Sheet: {sheet_name} ===")
                
                # Extract data from each row
                for row in worksheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        row_text = '\t'.join(str(cell) if cell is not None else '' for cell in row)
                        text_content.append(row_text)
                
                text_content.append('')  # Empty line between sheets
            
            full_text = '\n'.join(text_content)
            
            # Create PDF from text using PyMuPDF
            pdf_path = file_path.replace('.xlsx', '.pdf')
            pdf_doc = fitz.open()  # Create new PDF
            page = pdf_doc.new_page()
            
            # Insert text with basic formatting
            text_rect = fitz.Rect(50, 50, 550, 750)
            page.insert_textbox(text_rect, full_text, fontsize=10, fontname="helv")
            
            pdf_doc.save(pdf_path)
            pdf_doc.close()
            
            return pdf_path
            
        except ImportError:
            logger.error("openpyxl not installed. Please install with: pip install openpyxl")
            raise ValueError("Excel processing requires openpyxl library")
        except Exception as e:
            logger.error(f"Error converting XLSX to PDF: {str(e)}")
            raise
    
    async def _convert_xls_to_pdf(self, file_path: str) -> str:
        """Convert XLS to PDF using xlrd"""
        try:
            import xlrd
            
            # Load the Excel workbook
            workbook = xlrd.open_workbook(file_path)
            text_content = []
            
            for sheet_name in workbook.sheet_names():
                worksheet = workbook.sheet_by_name(sheet_name)
                text_content.append(f"=== Sheet: {sheet_name} ===")
                
                # Extract data from each row
                for row_idx in range(worksheet.nrows):
                    row_values = []
                    for col_idx in range(worksheet.ncols):
                        cell = worksheet.cell(row_idx, col_idx)
                        row_values.append(str(cell.value) if cell.value else '')
                    
                    if any(val for val in row_values):
                        row_text = '\t'.join(row_values)
                        text_content.append(row_text)
                
                text_content.append('')  # Empty line between sheets
            
            full_text = '\n'.join(text_content)
            
            # Create PDF from text using PyMuPDF
            pdf_path = file_path.replace('.xls', '.pdf')
            pdf_doc = fitz.open()  # Create new PDF
            page = pdf_doc.new_page()
            
            # Insert text with basic formatting
            text_rect = fitz.Rect(50, 50, 550, 750)
            page.insert_textbox(text_rect, full_text, fontsize=10, fontname="helv")
            
            pdf_doc.save(pdf_path)
            pdf_doc.close()
            
            return pdf_path
            
        except ImportError:
            logger.error("xlrd not installed. Please install with: pip install xlrd")
            raise ValueError("Legacy Excel processing requires xlrd library")
        except Exception as e:
            logger.error(f"Error converting XLS to PDF: {str(e)}")
            raise
    
    async def _convert_pptx_to_pdf(self, file_path: str) -> str:
        """Convert PPTX to PDF using python-pptx"""
        try:
            from pptx import Presentation
            
            # Load the PowerPoint presentation
            presentation = Presentation(file_path)
            text_content = []
            
            for slide_idx, slide in enumerate(presentation.slides, 1):
                text_content.append(f"=== Slide {slide_idx} ===")
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text)
                
                text_content.append('')  # Empty line between slides
            
            full_text = '\n'.join(text_content)
            
            # Create PDF from text using PyMuPDF
            pdf_path = file_path.replace('.pptx', '.pdf')
            pdf_doc = fitz.open()  # Create new PDF
            page = pdf_doc.new_page()
            
            # Insert text with basic formatting
            text_rect = fitz.Rect(50, 50, 550, 750)
            page.insert_textbox(text_rect, full_text, fontsize=11, fontname="helv")
            
            pdf_doc.save(pdf_path)
            pdf_doc.close()
            
            return pdf_path
            
        except ImportError:
            logger.error("python-pptx not installed. Please install with: pip install python-pptx")
            raise ValueError("PowerPoint processing requires python-pptx library")
        except Exception as e:
            logger.error(f"Error converting PPTX to PDF: {str(e)}")
            raise
    
    async def _convert_ppt_to_pdf(self, file_path: str) -> str:
        """Convert PPT to PDF (basic implementation)"""
        # For legacy .ppt files, we'll create a simple PDF with a conversion notice
        # In production, you might want to use a more sophisticated conversion tool
        pdf_path = file_path.replace('.ppt', '.pdf')
        
        # Create a simple PDF with a message about PPT conversion
        pdf_doc = fitz.open()
        page = pdf_doc.new_page()
        
        text = "Legacy PowerPoint (PPT) file uploaded. For better text extraction, please upload PPTX or PDF format."
        text_rect = fitz.Rect(50, 50, 550, 750)
        page.insert_textbox(text_rect, text, fontsize=11, fontname="helv")
        
        pdf_doc.save(pdf_path)
        pdf_doc.close()
        
        return pdf_path
    
    async def _convert_image_to_pdf(self, file_path: str) -> str:
        """Convert image to PDF using PyMuPDF"""
        try:
            # Open image and convert to PDF
            source_image = Image.open(file_path)
            
            # Convert to RGB if necessary
            if source_image.mode != 'RGB':
                source_image = source_image.convert('RGB')
            
            # Create PDF from image
            pdf_path = os.path.splitext(file_path)[0] + '.pdf'
            pdf_doc = fitz.open()
            
            # Save image to bytes in PNG format for PyMuPDF
            img_buffer = io.BytesIO()
            source_image.save(img_buffer, format='PNG')
            img_bytes = img_buffer.getvalue()
            
            img_width, img_height = source_image.size
            
            # Create page with appropriate size
            page = pdf_doc.new_page(width=img_width, height=img_height)
            
            # Insert image into PDF using correct PyMuPDF method
            img_rect = fitz.Rect(0, 0, img_width, img_height)
            page.insert_image(img_rect, stream=img_bytes)
            
            pdf_doc.save(pdf_path)
            pdf_doc.close()
            
            return pdf_path
            
        except Exception as e:
            logger.error(f"Error converting image to PDF: {str(e)}")
            raise
    
    async def _convert_csv_to_pdf(self, file_path: str) -> str:
        """Convert CSV to PDF using pandas and PyMuPDF"""
        try:
            import pandas as pd
            
            # Read CSV with error handling
            try:
                df = pd.read_csv(file_path, encoding='utf-8')
            except UnicodeDecodeError:
                # Try different encodings
                try:
                    df = pd.read_csv(file_path, encoding='latin-1')
                except:
                    df = pd.read_csv(file_path, encoding='cp1252')
            
            # Create readable text format
            csv_content = []
            csv_content.append(f"CSV File: {os.path.basename(file_path)}")
            csv_content.append(f"Dimensions: {len(df)} rows × {len(df.columns)} columns")
            csv_content.append("")
            csv_content.append("Column Headers: " + " | ".join(str(col) for col in df.columns))
            csv_content.append("=" * 80)
            
            # Add data rows (limit to 1000 for performance)
            max_rows = 1000
            for idx, row in df.head(max_rows).iterrows():
                row_text = " | ".join(str(val) if pd.notna(val) else "" for val in row.values)
                csv_content.append(f"Row {idx + 1}: {row_text}")
            
            if len(df) > max_rows:
                csv_content.append(f"... ({len(df) - max_rows} additional rows truncated)")
            
            full_text = '\n'.join(csv_content)
            
            # Convert to PDF using PyMuPDF
            pdf_path = file_path.replace('.csv', '.pdf')
            pdf_doc = fitz.open()
            page = pdf_doc.new_page()
            
            # Insert text with formatting
            text_rect = fitz.Rect(50, 50, 550, 750)
            page.insert_textbox(text_rect, full_text, fontsize=8, fontname="helv")
            
            pdf_doc.save(pdf_path)
            pdf_doc.close()
            
            logger.info(f"CSV converted to PDF: {pdf_path}")
            return pdf_path
            
        except ImportError:
            logger.error("pandas required for CSV processing. Install with: pip install pandas")
            raise ValueError("CSV processing requires pandas library")
        except Exception as e:
            logger.error(f"CSV conversion failed: {str(e)}")
            raise
    
    async def _extract_text_with_ocr(self, pdf_path: str) -> str:
        """Extract text from PDF using PyMuPDF and OCR fallback"""
        try:
            extracted_text = []
            
            # Open PDF with PyMuPDF
            pdf_doc = fitz.open(pdf_path)
            
            for page_num in range(pdf_doc.page_count):
                page = pdf_doc[page_num]
                
                # First, try to extract text directly
                page_text = page.get_text()
                
                if page_text.strip():
                    # Text found, use it
                    sanitized_page_text = self._sanitize_extracted_text(page_text)
                    extracted_text.append(sanitized_page_text)
                else:
                    # No text found, use OCR
                    logger.info(f"No text found on page {page_num + 1}, using OCR")
                    
                    # Get page as image
                    page_pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for better OCR
                    img_data = page_pixmap.tobytes("png")
                    
                    # Convert to PIL Image
                    page_image = Image.open(io.BytesIO(img_data))
                    
                    # Use Tesseract OCR
                    ocr_text = pytesseract.image_to_string(page_image, lang='eng')
                    
                    if ocr_text.strip():
                        sanitized_ocr_text = self._sanitize_extracted_text(ocr_text)
                        extracted_text.append(sanitized_ocr_text)
                    
                    page_pixmap = None  # Free memory
            
            pdf_doc.close()
            
            # Combine all extracted text
            full_text = '\n\n'.join(extracted_text)
            
            if not full_text.strip():
                return "No text could be extracted from this document."
            
            # Sanitize the text before returning
            sanitized_text = self._sanitize_extracted_text(full_text.strip())
            logger.info(f"Text extraction complete. Original: {len(full_text)} chars, Sanitized: {len(sanitized_text)} chars")
            
            return sanitized_text
            
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {str(e)}")
            return f"Error extracting text: {str(e)}"
    
    async def process_document_with_images(self, file_path: str, original_filename: str) -> Tuple[str, str, str, Dict]:
        """
        Enhanced processing that preserves both text and images
        
        Args:
            file_path: Path to the uploaded file
            original_filename: Original filename with extension
            
        Returns:
            Tuple of (document_id, package_path, extracted_text, metadata)
        """
        start_time = time.time()
        document_id = str(uuid.uuid4())
        
        try:
            # Get file extension
            file_ext = Path(original_filename).suffix.lower().lstrip('.')
            
            if file_ext not in self.supported_formats:
                raise ValueError(f"Unsupported file format: {file_ext}")
            
            # Convert to PDF if needed
            if file_ext == 'pdf':
                pdf_path = file_path
            else:
                pdf_path = await self.supported_formats[file_ext](file_path)
            
            # Create document package directory
            package_dir = Path(f"storage/document_packages/{document_id}")
            package_dir.mkdir(parents=True, exist_ok=True)
            
            # Extract both text and images
            extracted_text, page_images, page_texts = await self._extract_text_and_images(pdf_path, package_dir)
            
            processing_time = time.time() - start_time
            
            # Create document metadata
            metadata = {
                "document_id": document_id,
                "original_filename": original_filename,
                "file_type": file_ext,
                "total_pages": len(page_images),
                "page_images": page_images,
                "page_texts": page_texts,
                "processing_timestamp": datetime.now().isoformat(),
                "processing_time": processing_time,
                "file_size": os.path.getsize(file_path)
            }
            
            # Save metadata
            with open(package_dir / "metadata.json", 'w', encoding='utf-8') as f:
                json.dump(metadata, f, indent=2)
            
            # Save combined text
            with open(package_dir / "combined_text.txt", 'w', encoding='utf-8') as f:
                f.write(extracted_text)
            
            logger.info(f"Document package created: {document_id}")
            logger.info(f"Pages processed: {len(page_images)}")
            logger.info(f"Text extracted: {len(extracted_text)} characters")
            
            return document_id, str(package_dir), extracted_text, metadata
            
        except Exception as e:
            logger.error(f"Error processing document with images {original_filename}: {str(e)}")
            raise
    
    async def _extract_text_and_images(self, pdf_path: str, package_dir: Path) -> Tuple[str, List[str], List[str]]:
        """Extract both text and images from PDF pages"""
        try:
            extracted_text = []
            page_images = []
            page_texts = []
            
            # Open PDF with PyMuPDF
            pdf_doc = fitz.open(pdf_path)
            
            # Create subdirectories
            (package_dir / "pages").mkdir(exist_ok=True)
            (package_dir / "text").mkdir(exist_ok=True)
            (package_dir / "thumbnails").mkdir(exist_ok=True)
            
            for page_num in range(pdf_doc.page_count):
                page = pdf_doc[page_num]
                
                # Extract text from page
                page_text = page.get_text()
                
                if page_text.strip():
                    # Text found directly, use it
                    page_text = self._sanitize_extracted_text(page_text)
                    extracted_text.append(f"--- Page {page_num + 1} ---\n{page_text}")
                    logger.info(f"Direct text extraction successful for page {page_num + 1}")
                else:
                    # No direct text, use OCR on the image
                    logger.info(f"Using OCR for page {page_num + 1}")
                    page_text = await self._ocr_page_from_pdf(page, page_num + 1)
                    page_text = self._sanitize_extracted_text(page_text)
                    extracted_text.append(f"--- Page {page_num + 1} (OCR) ---\n{page_text}")
                
                # Extract high-quality page image
                img_path = await self._extract_page_image(page, page_num, package_dir)
                page_images.append(str(img_path))
                
                # Save individual page text
                text_path = package_dir / "text" / f"page_{page_num + 1}.txt"
                with open(text_path, 'w', encoding='utf-8') as f:
                    f.write(page_text)
                page_texts.append(str(text_path))
                
                # Generate thumbnail
                await self._generate_thumbnail(img_path, package_dir / "thumbnails" / f"thumb_{page_num + 1}.png")
            
            pdf_doc.close()
            
            # Combine all extracted text
            full_text = '\n\n'.join(extracted_text)
            
            if not full_text.strip():
                full_text = "No text could be extracted from this document."
            else:
                # Sanitize the combined text
                full_text = self._sanitize_extracted_text(full_text.strip())
            
            return full_text, page_images, page_texts
            
        except Exception as e:
            logger.error(f"Error extracting text and images from PDF: {str(e)}")
            raise
    
    async def _extract_page_image(self, page, page_num: int, package_dir: Path) -> Path:
        """Extract a single page as high-quality image"""
        try:
            # Use high resolution for better quality
            zoom_matrix = fitz.Matrix(3.0, 3.0)  # 3x zoom for high quality
            page_pixmap = page.get_pixmap(matrix=zoom_matrix)
            
            # Save page image
            img_path = package_dir / "pages" / f"page_{page_num + 1}.png"
            page_pixmap.save(str(img_path))
            
            logger.info(f"Page {page_num + 1} image saved: {img_path}")
            return img_path
            
        except Exception as e:
            logger.error(f"Error extracting page {page_num + 1} image: {str(e)}")
            raise
    
    async def _ocr_page_from_pdf(self, page, page_num: int) -> str:
        """Extract text from a PDF page using OCR"""
        try:
            # Convert page to image with higher resolution for better OCR
            ocr_zoom_matrix = fitz.Matrix(2.0, 2.0)  # 2x zoom for OCR
            page_pixmap = page.get_pixmap(matrix=ocr_zoom_matrix)
            
            # Convert to PIL Image
            img_data = page_pixmap.tobytes("png")
            page_image = Image.open(io.BytesIO(img_data))
            
            # Perform OCR
            ocr_text = pytesseract.image_to_string(
                page_image,
                lang='eng',
                config='--psm 3 --oem 3'
            )
            
            # Clean up
            page_pixmap = None
            
            # Sanitize OCR text before returning
            sanitized_text = self._sanitize_extracted_text(ocr_text.strip())
            return sanitized_text
            
        except Exception as e:
            logger.error(f"Error in OCR for page {page_num}: {str(e)}")
            return f"[OCR Error on page {page_num}: {str(e)}]"
    
    async def _generate_thumbnail(self, img_path: Path, thumb_path: Path) -> None:
        """Generate thumbnail from page image"""
        try:
            with Image.open(img_path) as img:
                # Create thumbnail (max 300x300)
                img.thumbnail((300, 300), Image.Resampling.LANCZOS)
                img.save(thumb_path, "PNG")
                logger.info(f"Thumbnail generated: {thumb_path}")
                
        except Exception as e:
            logger.error(f"Error generating thumbnail: {str(e)}")
    
    def get_document_package_info(self, document_id: str) -> Optional[Dict]:
        """Get information about a document package"""
        try:
            metadata_path = Path(f"storage/document_packages/{document_id}/metadata.json")
            if metadata_path.exists():
                with open(metadata_path, 'r', encoding='utf-8') as f:
                    return json.load(f)
            return None
        except Exception as e:
            logger.error(f"Error reading document package info: {str(e)}")
            return None
    
    # ========== MULTIMODAL CLIP PROCESSING ==========
    
    async def process_document_multimodal(self, file_path: str, original_filename: str, max_images: int = 10) -> Tuple[str, str, float, Dict]:
        """
        Process document with both text extraction and CLIP image analysis
        
        Args:
            file_path: Path to the uploaded file
            original_filename: Original filename with extension
            max_images: Maximum number of images to process with CLIP
            
        Returns:
            Tuple of (pdf_path, extracted_text, processing_time, clip_data)
        """
        start_time = time.time()
        
        try:
            logger.info(f"Starting multimodal processing for: {original_filename}")
            
            # Step 1: Regular document processing (text extraction)
            pdf_path, extracted_text, base_processing_time = await self.process_document(file_path, original_filename)
            
            # Step 2: CLIP image processing
            clip_data = {'images': [], 'num_images': 0, 'clip_success': False}
            
            try:
                from services.clip_service import get_clip_service
                clip_service = get_clip_service()
                
                if clip_service:
                    logger.info(f"Processing images with CLIP for: {original_filename}")
                    
                    # Process document with CLIP
                    clip_result = await clip_service.process_document_images(
                        pdf_path=pdf_path,
                        extracted_text=extracted_text,
                        max_images=max_images
                    )
                    
                    if clip_result['success']:
                        clip_data = {
                            'images': clip_result['images'],
                            'num_images': clip_result['num_images'], 
                            'clip_success': True,
                            'message': clip_result['message']
                        }
                        logger.info(f"CLIP processing successful: {clip_result['num_images']} images analyzed")
                    else:
                        logger.warning(f"CLIP processing failed: {clip_result.get('error', 'Unknown error')}")
                        clip_data['error'] = clip_result.get('error', 'CLIP processing failed')
                else:
                    logger.warning("CLIP service not available")
                    clip_data['error'] = 'CLIP service not available'
                    
            except Exception as clip_error:
                logger.error(f"Error in CLIP processing for {original_filename}: {str(clip_error)}")
                clip_data['error'] = str(clip_error)
            
            # Calculate total processing time
            total_processing_time = time.time() - start_time
            
            logger.info(f"Multimodal processing completed for {original_filename}")
            logger.info(f"Text processing: {base_processing_time:.2f}s, Total: {total_processing_time:.2f}s")
            logger.info(f"Images processed: {clip_data['num_images']}")
            
            return pdf_path, extracted_text, total_processing_time, clip_data
            
        except Exception as e:
            processing_time = time.time() - start_time
            logger.error(f"Error in multimodal document processing for {original_filename}: {str(e)}")
            
            # Return with error in clip_data
            error_clip_data = {
                'images': [],
                'num_images': 0,
                'clip_success': False,
                'error': str(e),
                'message': 'Multimodal processing failed'
            }
            
            # Try to return at least the text processing result
            try:
                pdf_path, extracted_text, _ = await self.process_document(file_path, original_filename)
                return pdf_path, extracted_text, processing_time, error_clip_data
            except:
                # Complete failure - return empty results
                return file_path, f"Error processing document: {str(e)}", processing_time, error_clip_data